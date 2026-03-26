"""
Add real screenshots and images to Alfalah_AI_2026_V3_FINAL.pptx
- Slide 5 (Azure Architecture): architecture_diagram.png
- Slide 9 (Live Demo): platform screenshot
- Slide 1 (Title): small platform screenshot inset
Run: python scripts/add_images_pptx.py
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE = Path(__file__).parent.parent
ARCH_IMG   = BASE / "architecture_diagram.png"
SCREEN_IMG = BASE / "docs" / "screenshots" / "screencapture-govrag-v3-func-azurewebsites-net-2026-03-26-16_24_26.png"

prs = Presentation(BASE / "Alfalah_AI_2026_V3_FINAL.pptx")

slide_w = prs.slide_width.inches   # 13.33"
slide_h = prs.slide_height.inches  # 7.5"
print(f"Slide size: {slide_w:.1f}\" x {slide_h:.1f}\"")

# ── SLIDE 5 — Azure Architecture: add architecture_diagram.png ────────────────
slide5 = prs.slides[4]
if ARCH_IMG.exists():
    # Place in lower-right area, half width
    pic = slide5.shapes.add_picture(
        str(ARCH_IMG),
        left=Inches(6.8),
        top=Inches(1.8),
        width=Inches(6.3),
        height=Inches(5.2),
    )
    print("Slide 5: architecture_diagram.png added")
else:
    print(f"WARN: {ARCH_IMG} not found")

# ── SLIDE 9 — Live Demo: add platform screenshot ──────────────────────────────
slide9 = prs.slides[8]
if SCREEN_IMG.exists():
    # Large screenshot, centered below the step list
    pic = slide9.shapes.add_picture(
        str(SCREEN_IMG),
        left=Inches(7.2),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(5.5),
    )
    print("Slide 9: platform screenshot added")
else:
    print(f"WARN: {SCREEN_IMG} not found")

# ── SLIDE 1 — Title: small platform screenshot inset ─────────────────────────
slide1 = prs.slides[0]
if SCREEN_IMG.exists():
    pic = slide1.shapes.add_picture(
        str(SCREEN_IMG),
        left=Inches(8.5),
        top=Inches(0.3),
        width=Inches(4.6),
        height=Inches(3.2),
    )
    print("Slide 1: platform screenshot inset added")

# ── SLIDE 7 — RAG Knowledge Engine: add architecture diagram as visual ────────
slide7 = prs.slides[6]
if ARCH_IMG.exists():
    pic = slide7.shapes.add_picture(
        str(ARCH_IMG),
        left=Inches(7.5),
        top=Inches(2.0),
        width=Inches(5.5),
        height=Inches(4.8),
    )
    print("Slide 7: architecture diagram added as RAG visual")

out = BASE / "Alfalah_AI_2026_V3_SUBMIT.pptx"
prs.save(str(out))
print(f"\nSaved: {out}")
print("Done.")
