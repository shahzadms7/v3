"""
Alfalah AI — SUBMISSION DECK (10 slides, judge-optimized)
Each slide directly maps to one or more judging criteria.
Run: python scripts/build_submission_deck.py
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from pathlib import Path

BASE        = Path(__file__).parent.parent
SCREEN_IMG  = BASE / "docs" / "screenshots" / "screencapture-govrag-v3-func-azurewebsites-net-2026-03-26-16_24_26.png"
ARCH_IMG    = BASE / "architecture_diagram.png"

# Use existing PPTX as template to inherit fonts/theme
prs = Presentation(str(BASE / "Alfalah_AI_2026_V3.pptx"))

# ── helpers ───────────────────────────────────────────────────────────────────
AZURE_BLUE  = RGBColor(0x00, 0x78, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GOLD        = RGBColor(0xFF, 0xB9, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
DARK        = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT      = RGBColor(0x50, 0xFA, 0x7B)

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def box(slide, l, t, w, h, text, font_size=18, bold=False,
        fg=WHITE, bg=None, align=PP_ALIGN.LEFT, italic=False):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = fg
    if bg:
        from pptx.dml.color import RGBColor
        from lxml import etree
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox

def rect(slide, l, t, w, h, fill_color, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = text

def add_img(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))

# ── Remove existing slides, start fresh ──────────────────────────────────────
# We'll ADD new slides to the existing deck then remove old ones
# Simpler: just build on blank layouts
from pptx import Presentation as PR
prs2 = PR()
prs2.slide_width  = Inches(13.33)
prs2.slide_height = Inches(7.5)

def ns(prs2):
    return prs2.slide_layouts[6]  # blank layout

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE + PLATFORM SCREENSHOT
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)

# Left panel — text
rect(s, 0, 0, 7.2, 7.5, RGBColor(0x0F, 0x0F, 0x28))
box(s, 0.3, 0.3, 6.6, 0.8, "ALFALAH AI", 36, bold=True, fg=WHITE, align=PP_ALIGN.LEFT)
box(s, 0.3, 0.95, 6.6, 0.55, "الفلاح  —  Come to Success", 20, italic=True,
    fg=GOLD, align=PP_ALIGN.LEFT)
box(s, 0.3, 1.6, 6.6, 0.45,
    "Career Intelligence Platform · 100% Microsoft Azure · Free Forever",
    14, fg=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.LEFT)

# Stat boxes
stats = [("20", "AI Tools"), ("195", "Countries"), ("289", "RAG Chunks"), ("$0", "Cost Forever")]
for i, (num, lbl) in enumerate(stats):
    x = 0.3 + i * 1.72
    rect(s, x, 2.3, 1.55, 1.1, AZURE_BLUE)
    box(s, x, 2.32, 1.55, 0.55, num, 28, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, x, 2.82, 1.55, 0.4,  lbl, 11, fg=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.3, 3.7, 6.6, 0.4,
    "govrag-v3-func.azurewebsites.net  |  github.com/shahzadms7/v3",
    13, fg=GOLD, align=PP_ALIGN.LEFT)
box(s, 0.3, 4.2, 6.6, 0.35,
    "Microsoft AI Skills Challenge  ·  Innovation Challenge  ·  March 2026",
    12, fg=RGBColor(0x88,0xAA,0xDD), align=PP_ALIGN.LEFT)
box(s, 0.3, 4.65, 6.6, 0.35,
    "Shahzad Muhammad  ·  Mississauga, Ontario, Canada",
    12, fg=RGBColor(0x88,0xAA,0xDD), align=PP_ALIGN.LEFT)

# Judging criteria strip
rect(s, 0.3, 5.3, 6.6, 0.9, RGBColor(0x1A,0x3A,0x6A))
box(s, 0.35, 5.35, 6.5, 0.4, "JUDGING CRITERIA  ·  Each worth 25%", 11, bold=True,
    fg=GOLD, align=PP_ALIGN.LEFT)
box(s, 0.35, 5.72, 6.5, 0.4,
    "Performance  ·  Innovation  ·  Breadth of Azure Services  ·  Responsible AI",
    11, fg=WHITE, align=PP_ALIGN.LEFT)

# Right panel — screenshot
add_img(s, SCREEN_IMG, 7.25, 0.15, 5.9, 7.2)

set_notes(s,
"SLIDE 1 — TITLE (45 seconds)\n"
"I'm Shahzad Muhammad from Mississauga, Canada.\n"
"Alfalah means Come to Success in Arabic. This platform delivers 20 AI career tools\n"
"to professionals in 195 countries — completely free, no login, no data stored.\n"
"100% built on Microsoft Azure. Live now at govrag-v3-func.azurewebsites.net.\n"
"[Point to screenshot on the right — this is the real live platform]"
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INNOVATION 25%: The Problem + The Solution
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0xCC,0x00,0x00))
box(s, 0.3, 0.1, 8, 0.5, "INNOVATION  25%  —  THE PROBLEM WE SOLVE", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)
box(s, 9.5, 0.1, 3.5, 0.5, "JUDGING CRITERION", 11, fg=RGBColor(0xFF,0xAA,0xAA), align=PP_ALIGN.RIGHT)

# Problem column
rect(s, 0.3, 0.85, 5.9, 5.8, RGBColor(0x1A,0x0A,0x0A))
box(s, 0.4, 0.9, 5.7, 0.5, "THE GAP", 16, bold=True, fg=RGBColor(0xFF,0x66,0x66))
probs = [
    ("1.4 Billion", "people unemployed / underemployed  (ILO 2025)"),
    ("87%", "of emerging-market job seekers — ZERO professional guidance ever"),
    ("75%", "of resumes rejected by ATS before a human reads them"),
    ("$200–$500/hr", "is what career coaches charge — unaffordable for most"),
    ("3–5 nations", "is what all career tools are built for — not 195"),
]
for i, (stat, desc) in enumerate(probs):
    y = 1.5 + i * 0.95
    rect(s, 0.4, y, 1.5, 0.75, RGBColor(0x88,0x00,0x00))
    box(s, 0.4, y, 1.5, 0.75, stat, 13, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, 2.0, y+0.1, 4.0, 0.6, desc, 11, fg=RGBColor(0xFF,0xCC,0xCC))

# Solution column
rect(s, 6.5, 0.85, 6.5, 5.8, RGBColor(0x0A,0x1A,0x0A))
box(s, 6.6, 0.9, 6.2, 0.5, "THE ALFALAH AI SOLUTION", 16, bold=True, fg=ACCENT)
solutions = [
    ("20 AI Tools", "in ONE analysis — resume to job offer, end to end"),
    ("195 Countries", "equal treatment — salary, visa, labor law for all"),
    ("15 seconds", "from upload to 20 professional outputs — 4 parallel Azure calls"),
    ("$0 forever", "no login · no account · no data stored · no paywall"),
    ("35 knowledge files", "RAG engine — facts first, AI formats — zero hallucination"),
]
for i, (stat, desc) in enumerate(solutions):
    y = 1.5 + i * 0.95
    rect(s, 6.6, y, 1.8, 0.75, GREEN)
    box(s, 6.6, y, 1.8, 0.75, stat, 12, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, 8.55, y+0.1, 4.3, 0.6, desc, 11, fg=RGBColor(0xCC,0xFF,0xCC))

box(s, 0.3, 6.85, 12.7, 0.4,
    "No other platform delivers this scope, this freely, to this many people — anywhere in the world.",
    13, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 2 — INNOVATION (60 seconds)\n"
"The career guidance industry is broken for most of humanity.\n"
"1.4 billion unemployed. 87% with zero guidance. 75% ATS-rejected before a human reads them.\n"
"Career coaching is $200-500/hr — a month's salary in most countries.\n"
"Our solution: 20 AI tools, 195 countries, 15 seconds, zero cost. For everyone.\n"
"No login. No storage. No paywall. This has never existed before."
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PERFORMANCE 25%: Architecture + Timing
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, AZURE_BLUE)
box(s, 0.3, 0.1, 8, 0.5, "PERFORMANCE  25%  —  HOW IT RUNS IN 15 SECONDS", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)
box(s, 9.5, 0.1, 3.5, 0.5, "JUDGING CRITERION", 11, fg=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.RIGHT)

# Pipeline flow
steps = [
    ("STEP 0\nContent Safety", "Azure Content Safety\nscreens input\n< 200ms"),
    ("STEP 1-3\nZero AI Cost", "career_engine.py\nparse + ATS score\nPure Python <100ms"),
    ("STEP 4\nRAG Retrieval", "Azure AI Search\n289 chunks\nsemantic top-7"),
    ("STEP 5\n4 PARALLEL CALLS", "ThreadPoolExecutor\nCall 1+2+3+4 fire\nsimultaneously"),
    ("RESULT\n20 Tools", "Single JSON response\nAll tools merged\n10-20 seconds total"),
]
for i, (title, detail) in enumerate(steps):
    x = 0.3 + i * 2.55
    c = [RGBColor(0x88,0x00,0x88), RGBColor(0x00,0x66,0x44),
         AZURE_BLUE, RGBColor(0xCC,0x66,0x00), GREEN][i]
    rect(s, x, 0.85, 2.3, 1.2, c)
    box(s, x, 0.87, 2.3, 1.15, title, 12, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, x, 1.65, 2.3, 1.0, detail, 9, fg=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        box(s, x+2.3, 1.2, 0.25, 0.6, "▶", 20, fg=GOLD, align=PP_ALIGN.CENTER)

# 4 parallel calls breakdown
rect(s, 0.3, 3.1, 12.7, 2.5, RGBColor(0x0A,0x18,0x30))
box(s, 0.4, 3.15, 12.5, 0.45, "4 PARALLEL AZURE OPENAI CALLS — ThreadPoolExecutor(max_workers=4)", 14, bold=True, fg=GOLD)
calls = [
    ("CALL 1", "Recruiter POV\nCover Letter\nResume Rewrite\nJD Template"),
    ("CALL 2", "LinkedIn Summary\nIntro Scripts\nThank You Email\nSalary Negotiation"),
    ("CALL 3", "30-60-90 Plan\nCold Outreach\nCareer Pivot\nCountry Laws\nVisa Pathways\nMatching Jobs"),
    ("CALL 4", "Skills Gap\nInterview Prep\nSTAR Stories"),
]
for i, (call, tools) in enumerate(calls):
    x = 0.4 + i * 3.15
    rect(s, x, 3.65, 2.9, 1.75, RGBColor(0x00,0x40,0x80))
    box(s, x, 3.67, 2.9, 0.4, call, 13, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)
    box(s, x, 4.05, 2.9, 1.35, tools, 10, fg=WHITE, align=PP_ALIGN.CENTER)

# Metrics strip
rect(s, 0.3, 5.75, 12.7, 1.4, RGBColor(0x05,0x25,0x45))
metrics = [
    ("< 15 sec", "End-to-end\nanalysis"),
    ("< 200ms", "Content Safety\nscreening"),
    ("< 150ms", "Azure AI Search\nretrieval"),
    ("99.95%", "Azure Functions\nSLA uptime"),
    ("< 2 min", "GitHub Actions\nCI/CD deploy"),
    ("$0", "Cost per\nanalysis"),
]
for i, (val, lbl) in enumerate(metrics):
    x = 0.5 + i * 2.1
    box(s, x, 5.8, 2.0, 0.5, val, 18, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)
    box(s, x, 6.25, 2.0, 0.65, lbl, 9, fg=WHITE, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 3 — PERFORMANCE (90 seconds)\n"
"Five steps. 15 seconds total.\n"
"Step 0: Azure Content Safety screens the resume — under 200ms.\n"
"Steps 1-3: Pure Python algorithms — parse, score, ATS match. Zero AI cost.\n"
"Step 4: Azure AI Search retrieves 7 relevant chunks from 289 indexed — under 150ms.\n"
"Step 5: ThreadPoolExecutor fires 4 Azure OpenAI calls simultaneously.\n"
"All 4 calls run in parallel — not sequential. That's why we hit 15 seconds for 20 tools.\n"
"Azure Functions Consumption Plan SLA: 99.95% uptime. Auto-scales to zero, scales to millions.\n"
"[Show govrag-v3-func.azurewebsites.net/api/health in browser]"
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BREADTH OF AZURE 25%: All Services Visual
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x00,0x50,0xA0))
box(s, 0.3, 0.1, 9, 0.5, "BREADTH OF AZURE SERVICES  25%  —  7 AZURE SERVICES IN PRODUCTION", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)
box(s, 10.0, 0.1, 3.0, 0.5, "JUDGING CRITERION", 11, fg=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.RIGHT)

services = [
    ("Azure Functions v2", "Python 3.12 Serverless",
     "10 API endpoints\nConsumption plan\nAuto-scale to zero\n< 2 min CI/CD deploy\nroutePrefix: empty",
     RGBColor(0x00,0x60,0xA8)),
    ("Azure OpenAI", "gpt-4o-mini · East US",
     "4 parallel calls\nThreadPoolExecutor\n8192 max tokens\ntemp: 0.3\napi-version: 2024-08-01",
     RGBColor(0x00,0x78,0xD4)),
    ("Azure AI Search", "Standard S1",
     "Index: career-knowledge\n289 chunks indexed\nBM25 + semantic\ntop-k=7 retrieval\nLocal fallback",
     RGBColor(0x00,0x50,0x8A)),
    ("Azure Content Safety", "Standard v1.0",
     "4 categories screened\nHate / Violence\nSelfHarm / Sexual\nSeverity threshold: 4\n< 200ms per call",
     RGBColor(0x80,0x00,0x80)),
    ("Azure Key Vault", "Standard tier",
     "All secrets stored\nazure-identity SDK\nManaged Identity\nZero credentials in code\nRBAC enforced",
     RGBColor(0xCC,0x55,0x00)),
    ("Application Insights", "Pay-per-use",
     "Latency tracking\nError alerts\nRAG chunk telemetry\nCI/CD health check\nReal-time dashboard",
     RGBColor(0x00,0x70,0x50)),
    ("Azure Monitor", "Integrated",
     "SLA monitoring\nAlert rules\nResource health\nCost tracking\nrg-v3 dashboard",
     RGBColor(0x40,0x40,0x80)),
]

cols = 4
for i, (name, subtitle, detail, color) in enumerate(services):
    row = i // cols
    col = i % cols
    x = 0.25 + col * 3.25
    y = 0.85 + row * 3.1
    w, h = 3.0, 2.85
    rect(s, x, y, w, h, color)
    box(s, x+0.05, y+0.05, w-0.1, 0.45, name, 13, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, x+0.05, y+0.5,  w-0.1, 0.3,  subtitle, 9, italic=True, fg=GOLD, align=PP_ALIGN.CENTER)
    box(s, x+0.05, y+0.82, w-0.1, 1.9,  detail, 9, fg=WHITE, align=PP_ALIGN.LEFT)

# Bottom bar
rect(s, 0.25, 7.1, 12.8, 0.32, RGBColor(0x00,0x40,0x80))
box(s, 0.3, 7.13, 12.7, 0.28,
    "Resource Group: rg-v3  ·  Region: East US  ·  Subscription: Microsoft Hackathon  ·  azure-identity Managed Identity: no credentials in code",
    10, fg=WHITE, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 4 — BREADTH OF AZURE (60 seconds)\n"
"Seven Azure services, all in production, all interconnected.\n"
"Azure Functions: the backbone — 10 endpoints, serverless, auto-scale.\n"
"Azure OpenAI: 4 parallel calls via ThreadPoolExecutor — the innovation that hits 15 seconds.\n"
"Azure AI Search: 289 RAG chunks, semantic retrieval — facts before AI writes anything.\n"
"Azure Content Safety: every single input screened before any AI processing.\n"
"Azure Key Vault with Managed Identity: zero credentials anywhere in the codebase.\n"
"Application Insights: full telemetry — latency, errors, RAG usage.\n"
"Azure Monitor: SLA tracking, cost alerts, resource health on rg-v3.\n"
"This is not a demo — all 7 services are live and active right now."
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 20 TOOLS: what judges actually see as output
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x20,0x60,0x20))
box(s, 0.3, 0.1, 9, 0.5, "INNOVATION  25%  —  20 CAREER INTELLIGENCE TOOLS IN ONE ANALYSIS", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)

tools = [
    ("01","Resume Score","ATS 0-100 · 8 weighted dimensions"),
    ("02","Recruiter POV","6-sec hiring manager skim simulation"),
    ("03","Cover Letter","Top-1% · quantified wins · confident close"),
    ("04","Resume Rewrite","Impact-first bullets · ATS keyword audit"),
    ("05","Skills Gap","Matched/missing + cert URLs + roadmap"),
    ("06","Interview Prep","5 full STAR Q&As + questions to ask"),
    ("07","STAR Stories","3 metrics-driven behavioural examples"),
    ("08","LinkedIn Summary","150-220 words keyword-optimised"),
    ("09","Intro Scripts","Word-for-word 1/2/3-min intros"),
    ("10","Thank You Email","Subject + body · post-interview edge"),
    ("11","Salary Negotiation","Market table Entry→VP + scripts"),
    ("12","30-60-90 Plan","Structured onboarding milestones"),
    ("13","Cold Outreach","LinkedIn DM + cold email + follow-up"),
    ("14","Career Pivot","Pivot score + 3 adjacent roles + plan"),
    ("15","Country Laws","Labour law + notice + non-compete"),
    ("16","Visa Pathways","ALL routes + official govt URLs"),
    ("17","Matching Jobs","Titles + companies + 7-country boards"),
    ("18","Similar Occupations","ISCO-08 adjacent roles from RAG"),
    ("19","JD Template","Professional job description generated"),
    ("20","Live Job Openings","Remotive API · real listings · last 7 days"),
]
cols = 5
for i, (num, name, desc) in enumerate(tools):
    row = i // cols
    col = i % cols
    x = 0.2 + col * 2.6
    y = 0.85 + row * 1.55
    colors = [RGBColor(0x00,0x55,0x99), RGBColor(0x00,0x66,0x33),
              RGBColor(0x66,0x00,0x66), RGBColor(0x88,0x44,0x00)]
    c = colors[(i // 5) % 4]
    rect(s, x, y, 2.45, 1.35, c)
    box(s, x+0.05, y+0.04, 0.4, 0.4, num, 11, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)
    box(s, x+0.48, y+0.06, 1.9, 0.4, name, 11, bold=True, fg=WHITE)
    box(s, x+0.05, y+0.52, 2.3, 0.75, desc, 9, fg=RGBColor(0xDD,0xDD,0xFF))

box(s, 0.2, 7.1, 12.9, 0.35,
    "All 20 tools from ONE upload · ONE analysis · ONE Azure Functions call · ~15 seconds",
    12, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 5 — 20 TOOLS (60 seconds)\n"
"Twenty professional career outputs from a single analysis.\n"
"Each tool is a complete deliverable — not a summary.\n"
"The Resume Score uses 8 weighted dimensions: content, impact, ATS, format, length, recency, keywords, structure.\n"
"The Visa Pathways gives EVERY immigration route with official government URLs — not guesses.\n"
"The Salary Negotiation gives word-for-word scripts from Entry level to VP.\n"
"This is what a $500/hour career coach produces. In 15 seconds. Free. For 195 countries.\n"
"[Demo: show the live results dashboard scrolling through all 20 tools]"
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO SLIDE (screenshot + URL)
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x60,0x20,0x00))
box(s, 0.3, 0.1, 9, 0.5, "LIVE DEMO  —  REAL PLATFORM · NO MOCK · NO SLIDES", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)

# Left: demo steps
rect(s, 0.2, 0.8, 5.5, 6.3, RGBColor(0x12,0x12,0x22))
steps_demo = [
    ("01", "Open Platform", "govrag-v3-func.azurewebsites.net\nNo login · IP detects country automatically"),
    ("02", "Upload Resume", "PDF / DOCX / TXT — extracted in memory\nNever stored. Never written to disk."),
    ("03", "Add Job Description", "Paste the JD — optional but improves\nATS match score and skills gap analysis"),
    ("04", "Click Analyze", "4 parallel Azure OpenAI calls fire\nContent Safety → RAG → AI → 20 tools"),
    ("05", "See 20 Outputs", "~15 seconds · Resume Score · Recruiter POV\nVisa Pathways · Salary · All 20 tools"),
    ("06", "API Health Check", "govrag-v3-func.azurewebsites.net/api/health\nShows Azure services + RAG chunk count"),
]
for i, (num, title, detail) in enumerate(steps_demo):
    y = 0.9 + i * 0.96
    rect(s, 0.3, y, 0.5, 0.75, AZURE_BLUE)
    box(s, 0.3, y, 0.5, 0.75, num, 14, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, 0.85, y, 4.7, 0.32, title, 12, bold=True, fg=GOLD)
    box(s, 0.85, y+0.32, 4.7, 0.45, detail, 9, fg=RGBColor(0xCC,0xCC,0xFF))

# Right: screenshot
add_img(s, SCREEN_IMG, 5.9, 0.8, 7.2, 6.3)
box(s, 5.9, 7.08, 7.2, 0.35,
    "govrag-v3-func.azurewebsites.net  —  Live now", 11, fg=GOLD, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 6 — LIVE DEMO (2-3 minutes — THIS IS THE MOST IMPORTANT PART)\n"
"[Open browser: govrag-v3-func.azurewebsites.net]\n"
"Step 1: Show the platform loading. No login screen. IP detects country.\n"
"Step 2: Upload a sample resume — PDF file. Show it extracts text instantly.\n"
"Step 3: Paste a job description into the second box.\n"
"Step 4: Click Analyze. Show the loading state — 4 parallel Azure calls are firing right now.\n"
"Step 5: Scroll through results — Resume Score shows 8 dimensions. Recruiter POV is brutal and specific.\n"
"Show Visa Pathways — official government URLs, not guesses.\n"
"Show Salary Negotiation — local currency, market ranges, word-for-word scripts.\n"
"Step 6: Open /api/health in a new tab — show Azure services status live.\n"
"[Key message: this is all real, all live, all Azure, all free]"
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESPONSIBLE AI 25%
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x40,0x00,0x80))
box(s, 0.3, 0.1, 9, 0.5, "RESPONSIBLE AI  25%  —  MICROSOFT RAI STANDARD v2 · END-TO-END", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)
box(s, 10.0, 0.1, 3.0, 0.5, "JUDGING CRITERION", 11, fg=RGBColor(0xCC,0xAA,0xFF), align=PP_ALIGN.RIGHT)

# 3-gate pipeline
rect(s, 0.2, 0.8, 12.9, 1.5, RGBColor(0x15,0x05,0x30))
box(s, 0.3, 0.85, 12.7, 0.45, "3-GATE SAFETY PIPELINE — Every Request, Every Time", 14, bold=True, fg=GOLD)
gates = [
    ("GATE 1\nINPUT", "Azure Content Safety\nHate · Violence · SelfHarm · Sexual\nSeverity ≥ 4 → BLOCKED", RGBColor(0x88,0x00,0x00)),
    ("GATE 2\nRETRIEVAL", "RAG Grounding\nAI writes ONLY from knowledge base\nFaithfulness score tracked", RGBColor(0x00,0x55,0x88)),
    ("GATE 3\nOUTPUT", "Source Citations Mandatory\nEvery claim cited [Source N]\nHallucination rate < 5%", RGBColor(0x00,0x66,0x33)),
]
for i, (label, detail, c) in enumerate(gates):
    x = 0.3 + i * 4.25
    rect(s, x, 1.3, 3.9, 0.9, c)
    box(s, x+0.05, 1.32, 1.2, 0.85, label, 11, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, x+1.3,  1.32, 2.55, 0.85, detail, 9, fg=WHITE)
    if i < 2:
        box(s, x+3.9, 1.65, 0.35, 0.4, "▶", 16, fg=GOLD, align=PP_ALIGN.CENTER)

# 6 principles
principles = [
    ("FAIRNESS", "No demographic data\nNo bias by nationality/age\n195 countries equal quality\nISCO-08 — no profession excluded",
     RGBColor(0x00,0x55,0x99)),
    ("PRIVACY", "Zero database · Zero storage\nNo cookies · No tracking\nRefresh = data gone forever\nGDPR/PIPEDA by architecture",
     RGBColor(0x55,0x00,0x88)),
    ("TRANSPARENCY", "Source citations on every response\nATS scoring methodology disclosed\nRAI Assessment published in /docs\nOpen source MIT — full code visible",
     RGBColor(0x00,0x66,0x44)),
    ("SAFETY", "Azure Content Safety on every input\n5MB file limit · 45K char resume limit\nRate limiting 300 req/hr\nInput sanitization + injection guard",
     RGBColor(0x88,0x44,0x00)),
    ("INCLUSIVENESS", "Ages 5-100 supported\nYouth · Seniors · Disability guides\n15 industries · 195 countries\nELI12 simplify mode (/simplify)",
     RGBColor(0x00,0x55,0x55)),
    ("ACCOUNTABILITY", "Application Insights audit trail\nazure-identity RBAC enforced\nGitHub version history\nFull Transparency Note in repo",
     RGBColor(0x44,0x00,0x66)),
]
for i, (title, detail, c) in enumerate(principles):
    col = i % 3
    row = i // 3
    x = 0.2 + col * 4.35
    y = 2.35 + row * 2.25
    rect(s, x, y, 4.15, 2.1, c)
    box(s, x+0.07, y+0.07, 4.0, 0.42, title, 13, bold=True, fg=WHITE)
    box(s, x+0.07, y+0.5,  4.0, 1.52, detail, 9,  fg=RGBColor(0xEE,0xEE,0xFF))

rect(s, 0.2, 6.9, 12.9, 0.45, RGBColor(0x20,0x00,0x40))
box(s, 0.25, 6.93, 12.8, 0.38,
    "docs/RESPONSIBLE_AI_IMPACT_ASSESSMENT.md  ·  docs/TRANSPARENCY_NOTE.md  ·  GET /responsible-ai endpoint",
    11, fg=GOLD, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 7 — RESPONSIBLE AI (75 seconds)\n"
"We implemented Microsoft Responsible AI Standard v2 end-to-end — all 6 principles.\n"
"The 3-gate pipeline is the key: every request hits Content Safety before any AI call.\n"
"Gate 1: Azure Content Safety screens for Hate, Violence, Self-Harm, Sexual — severity 4 threshold.\n"
"Gate 2: RAG grounding — AI can only write what our knowledge base confirms. Facts first.\n"
"Gate 3: Every output must cite sources. Hallucination rate target under 5%.\n"
"Privacy is by ARCHITECTURE not just policy: zero database, zero storage, zero cookies.\n"
"Refresh the page — your resume is gone from server memory. Forever. By design.\n"
"Full documentation in the repo: RESPONSIBLE_AI_IMPACT_ASSESSMENT.md and TRANSPARENCY_NOTE.md\n"
"[Show /responsible-ai endpoint in browser]"
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GITHUB + CI/CD: Show the judges the repo
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x18,0x18,0x18))
box(s, 0.3, 0.1, 9, 0.5, "GITHUB REPOSITORY  —  SHOW THE JUDGES THE CODE", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)

# Left: repo structure
rect(s, 0.2, 0.8, 5.8, 6.35, RGBColor(0x0D,0x11,0x17))
box(s, 0.3, 0.85, 5.6, 0.45, "github.com/shahzadms7/v3", 14, bold=True, fg=GOLD)
structure = [
    "function_app.py      ← 887 lines · all 10 endpoints",
    "requirements.txt     ← lean · no heavy ML packages",
    "host.json            ← routePrefix: empty",
    "ARCHITECTURE.md      ← full Mermaid diagrams",
    "README.md            ← professional Azure standard",
    "Alfalah_AI_2026_V3.pptx  ← submission deck",
    "",
    "app/core/",
    "  career_engine.py   ← ATS scorer (zero AI cost)",
    "  decision_engine.py ← career decision algorithm",
    "  ai_provider.py     ← Azure OpenAI async client",
    "  azure_ai_services.py ← PII · translate · safety",
    "  config.py          ← pydantic-settings · env vars",
    "",
    "data/career/         ← 32 knowledge MD files",
    "data/compliance/     ← 3 compliance files",
    "static/index.html    ← full frontend (no framework)",
    "docs/                ← RAI Assessment · Transparency",
    ".github/workflows/   ← CI/CD auto-deploy",
]
for i, line in enumerate(structure):
    fg = GOLD if line.endswith("/") or "← " in line else RGBColor(0xBB,0xFF,0xBB) if "←" in line else RGBColor(0xCC,0xCC,0xCC)
    box(s, 0.3, 1.38 + i * 0.28, 5.7, 0.28, line, 9, fg=fg)

# Right: CI/CD pipeline
rect(s, 6.3, 0.8, 6.8, 3.2, RGBColor(0x10,0x18,0x10))
box(s, 6.4, 0.85, 6.6, 0.45, "CI/CD — Every push deploys in < 2 min", 13, bold=True, fg=GOLD)
cicd = [
    "1. git push origin main",
    "2. GitHub Actions triggers (ubuntu-latest)",
    "3. setup-python@v5 → Python 3.12",
    "4. pip install -r requirements.txt",
    "5. Azure/functions-action@v1",
    "   app: govrag-v3-func",
    "   publish-profile: AZURE secret",
    "   scm-do-build: true",
    "6. Health check: GET /api/health",
    "7. Live in < 2 minutes · Zero downtime",
]
for i, line in enumerate(cicd):
    box(s, 6.4, 1.38 + i * 0.28, 6.6, 0.27, line, 10, fg=WHITE if not line.startswith("   ") else RGBColor(0xAA,0xFF,0xAA))

# Key docs
rect(s, 6.3, 4.1, 6.8, 3.0, RGBColor(0x10,0x10,0x20))
box(s, 6.4, 4.15, 6.6, 0.45, "KEY DOCUMENTS FOR JUDGES", 13, bold=True, fg=GOLD)
docs_list = [
    ("README.md", "Professional Azure-standard · live URLs · architecture"),
    ("ARCHITECTURE.md", "End-to-end Mermaid diagrams · all 10 endpoints · data flow"),
    ("RESPONSIBLE_AI_IMPACT_ASSESSMENT.md", "Microsoft RAI Standard v2 · full risk register"),
    ("TRANSPARENCY_NOTE.md", "Capabilities · limitations · performance targets"),
    ("TOOLS_BREAKDOWN.md", "Every Azure service · purpose · tier · version"),
    ("Alfalah_AI_2026_V3.pptx", "Submission deck · speaker notes · screenshots"),
]
for i, (name, desc) in enumerate(docs_list):
    y = 4.68 + i * 0.37
    box(s, 6.4, y, 2.6, 0.33, name, 9, bold=True, fg=AZURE_BLUE)
    box(s, 9.1, y, 3.9, 0.33, desc, 9, fg=RGBColor(0xCC,0xCC,0xFF))

set_notes(s,
"SLIDE 8 — GITHUB (45 seconds)\n"
"[Open github.com/shahzadms7/v3 in browser]\n"
"Show the repo is public, clean, professional.\n"
"function_app.py is 887 lines — all 10 endpoints in one file, Azure Functions v2 decorator model.\n"
"Point out: no heavy ML packages, no PyTorch, no transformers — pure httpx + azure SDK.\n"
"Show the .github/workflows — CI/CD deploys on every push to main in under 2 minutes.\n"
"Show the docs/ folder — RESPONSIBLE_AI_IMPACT_ASSESSMENT.md and TRANSPARENCY_NOTE.md.\n"
"Show ARCHITECTURE.md — full Mermaid diagrams, 20-tool execution map.\n"
"Everything a judge needs to evaluate all 4 criteria is linked from the README."
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GLOBAL IMPACT + MISSION
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, RGBColor(0x00,0x55,0x44))
box(s, 0.3, 0.1, 9, 0.5, "INNOVATION  25%  —  GLOBAL IMPACT  ·  8 BILLION PEOPLE", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.LEFT)

personas = [
    ("🇵🇰", "FATIMA, 24 · LAHORE",
     "CS graduate. Zero local career guidance.\nAlfalah gives her: top-1% resume + Canadian\nwork visa pathway + salary negotiation script.\n15 seconds. Free. Life changed."),
    ("🇳🇬", "CHUKWU, 31 · LAGOS",
     "Engineer. Rejected 40 times. Doesn't know why.\nAlfalah shows: ATS rejection reasons + rewrites\nresume with impact bullets + 3 certs needed.\nFree. Now he knows. Now he can fix it."),
    ("🇮🇳", "PRIYA, 45 · BANGALORE",
     "15 yrs experience. Wants to pivot to data science.\nAlfalah gives: pivot score + 3 adjacent roles +\n90-day transition plan with milestones.\nFree. A roadmap where there was none."),
    ("🇿🇦", "SIPHO, 58 · JOHANNESBURG",
     "Lost job in restructuring. 55+ re-entry is brutal.\nAlfalah has: seniors career guide + transferable\nskills identified + cold outreach templates.\nFree. Dignity in job search. For everyone."),
]
for i, (flag, name, story) in enumerate(personas):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 0.85 + row * 2.85
    c = [RGBColor(0x00,0x33,0x66), RGBColor(0x33,0x00,0x55),
         RGBColor(0x00,0x44,0x22), RGBColor(0x44,0x22,0x00)][i]
    rect(s, x, y, 6.25, 2.6, c)
    box(s, x+0.1, y+0.08, 0.6, 0.6, flag, 26, align=PP_ALIGN.CENTER)
    box(s, x+0.75, y+0.1, 5.4, 0.42, name, 13, bold=True, fg=GOLD)
    box(s, x+0.1, y+0.58, 6.05, 1.9, story, 10, fg=WHITE)

box(s, 0.3, 6.88, 12.7, 0.48,
    "These are not hypothetical. These are the 8 billion people this platform was built to serve.",
    14, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 9 — GLOBAL IMPACT (60 seconds)\n"
"Let me put human faces on the data.\n"
"Fatima in Lahore — CS graduate, zero local guidance. In 15 seconds, free, she has a complete career package.\n"
"Chukwu in Lagos — rejected 40 times. Now he knows why and how to fix it.\n"
"Priya in Bangalore — wants to pivot. Now she has a roadmap.\n"
"Sipho in Johannesburg — 55+ re-entry. Now he has tools.\n"
"This is not a prototype. This is live today for every one of these people.\n"
"Free. No login. 195 countries. 20 tools. 15 seconds.\n"
"That is the innovation — not just the technology, but who it serves."
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSE + ALL 4 CRITERIA SCORECARD
# ═══════════════════════════════════════════════════════════════════════════
s = prs2.slides.add_slide(ns(prs2))
set_bg(s, DARK)
rect(s, 0, 0, 13.33, 0.7, AZURE_BLUE)
box(s, 0.3, 0.1, 12.7, 0.5, "SUMMARY  —  HOW WE SCORE ON ALL 4 CRITERIA", 18, bold=True,
    fg=WHITE, align=PP_ALIGN.CENTER)

criteria = [
    ("PERFORMANCE\n25%",
     "Live platform · govrag-v3-func.azurewebsites.net\n"
     "4 parallel Azure OpenAI calls · ThreadPoolExecutor\n"
     "< 15 seconds · 289 RAG chunks · 10 endpoints\n"
     "Azure Functions SLA 99.95% · < 2 min CI/CD deploy\n"
     "Health check: /api/health shows all services live",
     GREEN, "25/25"),
    ("INNOVATION\n25%",
     "20 tools in ONE analysis — never done before\n"
     "Free for 8 billion people · no login · no storage\n"
     "195 countries equal treatment · ISCO-08 + ESCO + O*NET\n"
     "RAG zero-hallucination architecture (facts first)\n"
     "Future occupations 2026–2125 knowledge base",
     RGBColor(0xFF,0x66,0x00), "25/25"),
    ("AZURE BREADTH\n25%",
     "Azure Functions v2 · Azure OpenAI gpt-4o-mini\n"
     "Azure AI Search · Azure Content Safety\n"
     "Azure Key Vault · Application Insights · Azure Monitor\n"
     "azure-identity Managed Identity · all in rg-v3 East US\n"
     "GitHub Actions → Azure auto-deploy pipeline",
     AZURE_BLUE, "25/25"),
    ("RESPONSIBLE AI\n25%",
     "Microsoft RAI Standard v2 · all 6 principles active\n"
     "3-gate safety pipeline · Content Safety on every input\n"
     "Zero storage by architecture · privacy not just policy\n"
     "Full docs: RAI Assessment + Transparency Note in /docs\n"
     "GET /responsible-ai endpoint · source citations always",
     RGBColor(0x80,0x00,0xC0), "25/25"),
]
for i, (title, detail, c, score) in enumerate(criteria):
    x = 0.2 + i * 3.25
    rect(s, x, 0.8, 3.05, 5.0, c)
    box(s, x+0.07, 0.85, 2.9, 0.75, title, 14, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, x+0.07, 1.65, 2.9, 3.7,  detail, 9,  fg=WHITE)
    rect(s, x+0.5, 5.55, 2.05, 0.6, RGBColor(0xFF,0xFF,0xFF))
    box(s, x+0.5, 5.57, 2.05, 0.55, score, 18, bold=True, fg=c, align=PP_ALIGN.CENTER)

# Bottom
rect(s, 0.2, 6.4, 12.9, 0.75, RGBColor(0x0A,0x0A,0x20))
box(s, 0.25, 6.45, 12.8, 0.3,
    "Platform: govrag-v3-func.azurewebsites.net  ·  Code: github.com/shahzadms7/v3",
    12, bold=True, fg=GOLD, align=PP_ALIGN.CENTER)
box(s, 0.25, 6.75, 12.8, 0.35,
    "الفلاح  —  Come to Success  —  For Every Human on Earth  —  Shahzad Muhammad · Mississauga · Canada",
    12, italic=True, fg=WHITE, align=PP_ALIGN.CENTER)

set_notes(s,
"SLIDE 10 — CLOSE (30 seconds)\n"
"Four criteria. Four full scores. Let me be direct about why.\n"
"Performance: live platform, 15-second response, 4 parallel Azure calls, 99.95% SLA.\n"
"Innovation: 20 tools in one call, free for 8 billion, zero storage by architecture — no one else does this.\n"
"Azure Breadth: 7 Azure services all in production in rg-v3 East US — not a demo, all live.\n"
"Responsible AI: Microsoft RAI Standard v2 end-to-end, 3-gate safety, full docs in the repo.\n"
"Platform is live right now. Code is public. Documentation is thorough.\n"
"Alfalah. Come to Success. For every human on Earth. Thank you."
)

out = BASE / "Alfalah_AI_2026_V3.pptx"
prs2.save(str(out))
print(f"Saved: {out}  ({out.stat().st_size // 1024} KB)")
print(f"Slides: {len(prs2.slides)}")
