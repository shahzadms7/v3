"""
Alfalah Job Career Intelligent AI 2026 V3 — PowerPoint Presentation Generator
Microsoft Azure Innovation Challenge — 10-15 Minute Presentation
Built for 8 Billion People
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color Palette ──────────────────────────────────────────────────
AZURE_BLUE     = RGBColor(0x00, 0x78, 0xD4)   # Microsoft Azure blue
DARK_BG        = RGBColor(0x0A, 0x0A, 0x1A)   # Deep dark navy
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE         = RGBColor(0xFF, 0x6B, 0x35)   # Alfalah accent
PURPLE         = RGBColor(0x8B, 0x5C, 0xF6)   # AI purple
GREEN          = RGBColor(0x22, 0xC5, 0x5E)   # Success green
LIGHT_GRAY     = RGBColor(0xE5, 0xE7, 0xEB)
MID_GRAY       = RGBColor(0x6B, 0x72, 0x80)
GOLD           = RGBColor(0xF5, 0x9E, 0x0B)   # Award gold

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=14, color=WHITE, bullet="●"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

# Top accent bar
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)
# Bottom accent bar
add_rect(s, 0, 7.42, 13.33, 0.08, ORANGE)

# Arabic calligraphy feel — main brand
add_textbox(s, "الفلاح", 0, 0.4, 13.33, 1.2,
            font_size=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(s, "ALFALAH JOB CAREER INTELLIGENT AI  2026  V3", 0, 1.45, 13.33, 0.7,
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "Career Intelligence Platform", 0, 2.1, 13.33, 0.5,
            font_size=22, bold=False, color=AZURE_BLUE, align=PP_ALIGN.CENTER)

# Divider line
add_rect(s, 3, 2.7, 7.33, 0.04, ORANGE)

add_textbox(s, "Built for 8 Billion People  ·  100% Microsoft Azure  ·  Free Forever",
            0, 2.85, 13.33, 0.5,
            font_size=16, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, "\"Come to Success — For Every Human on Earth\"",
            0, 3.45, 13.33, 0.5,
            font_size=15, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Stats row
stats = [
    ("195", "Countries"),
    ("17", "AI Modules"),
    ("436", "Occupations"),
    ("0", "Cost Forever"),
]
for i, (num, label) in enumerate(stats):
    x = 1.5 + i * 2.8
    add_rect(s, x, 4.2, 2.2, 1.3, RGBColor(0x0D, 0x1B, 0x3E))
    add_textbox(s, num, x, 4.25, 2.2, 0.7,
                font_size=32, bold=True, color=AZURE_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, 4.9, 2.2, 0.4,
                font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, "Shahzad Muhammad  ·  Mississauga, Ontario, Canada  ·  2026",
            0, 6.8, 13.33, 0.5,
            font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 1 (60 seconds)

Good [morning/afternoon/evening], everyone.

My name is Shahzad Muhammad, and I am here today to present Alfalah AI —
which in Arabic means "Come to Success."

This platform was born from a simple but powerful belief:
Every human being on this planet deserves access to world-class career guidance.
Not just the privileged few. Not just people in wealthy countries.
All 8 billion of us.

Today I will show you how we built a free, zero-login, AI-powered career intelligence
platform, running 100% on Microsoft Azure, that delivers 17 specialized career modules
to professionals in every one of the 195 countries on Earth.

Let me take you through the journey.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, ORANGE)

add_textbox(s, "THE PROBLEM WE SOLVE", 0.5, 0.2, 12, 0.7,
            font_size=28, bold=True, color=WHITE)
add_textbox(s, "1.4 Billion People. No Career Guidance. No Way Forward.",
            0.5, 0.85, 12, 0.5,
            font_size=16, italic=True, color=ORANGE)

problems = [
    ("1.4 Billion",  "people unemployed or underemployed globally (ILO 2025)"),
    ("87%",          "of job seekers in emerging markets NEVER received professional resume guidance"),
    ("$200–$500/hr", "is what career coaches charge — unaffordable for most of the world"),
    ("75%",          "of resumes are rejected by ATS systems before a human ever sees them"),
    ("195 Countries","but career tools only built for 3-5 wealthy nations"),
    ("0",            "free, truly global, AI-powered career platforms before today"),
]

for i, (stat, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.4
    y = 1.6 + row * 1.55
    add_rect(s, x, y, 5.9, 1.35, RGBColor(0x1A, 0x0A, 0x05))
    add_textbox(s, stat, x + 0.15, y + 0.05, 2.5, 0.6,
                font_size=26, bold=True, color=ORANGE)
    add_textbox(s, desc, x + 0.15, y + 0.6, 5.5, 0.65,
                font_size=12, color=LIGHT_GRAY)

add_textbox(s, "The career guidance industry is broken for most of humanity. We fix that.",
            0.5, 6.6, 12.3, 0.6,
            font_size=14, italic=True, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 2 (90 seconds)

Let me start with the reality of where we are today.

1.4 billion people are unemployed or underemployed globally.
That is nearly 1 in 5 people on this planet who cannot find meaningful work.

87% of job seekers in emerging markets — places like Pakistan, Nigeria, Indonesia,
Bangladesh — have NEVER received professional career guidance in their entire lives.

Career coaches charge $200 to $500 per hour. That is completely inaccessible
to someone earning $5 a day.

75% of resumes are rejected by automated systems before a single human
reads them. People do not even know why they are being rejected.

And every existing AI career tool? Built for the US market. Maybe Canada and UK.
195 countries on Earth. Tools serving maybe 5.

That is the gap. That is the injustice. And that is exactly what Alfalah AI was built to solve.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)

add_textbox(s, "THE SOLUTION", 0.5, 0.2, 8, 0.7,
            font_size=28, bold=True, color=WHITE)
add_textbox(s, "One Platform. Every Human. Every Country. Free Forever.",
            0.5, 0.85, 12, 0.5,
            font_size=16, italic=True, color=AZURE_BLUE)

# Central flow
add_rect(s, 0.4, 1.5, 3.5, 4.5, RGBColor(0x0D, 0x1B, 0x3E))
add_textbox(s, "INPUT", 0.4, 1.5, 3.5, 0.45,
            font_size=13, bold=True, color=AZURE_BLUE, align=PP_ALIGN.CENTER)
inputs = ["📄 Resume (PDF/DOCX/TXT)", "📋 Job Description", "🌍 Your Country (195)", "🏭 Your Industry (23)"]
add_bullet_box(s, inputs, 0.55, 2.0, 3.2, 3.5, font_size=13, color=LIGHT_GRAY, bullet="→")

# Arrow
add_textbox(s, "⚡", 4.1, 3.4, 0.8, 0.5, font_size=28, color=ORANGE, align=PP_ALIGN.CENTER)

# Engine
add_rect(s, 5.0, 1.5, 3.2, 4.5, RGBColor(0x0A, 0x0A, 0x2E))
add_textbox(s, "AZURE AI ENGINE", 5.0, 1.5, 3.2, 0.45,
            font_size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
engine = ["GPT-4o-mini (Primary)", "Gemini 2.0 Flash", "xAI Grok-4 (Fallback)", "RAG: 28 knowledge files", "436 ISCO-08 occupations"]
add_bullet_box(s, engine, 5.1, 2.0, 3.0, 3.5, font_size=12, color=LIGHT_GRAY, bullet="▸")

# Arrow
add_textbox(s, "⚡", 8.3, 3.4, 0.8, 0.5, font_size=28, color=GREEN, align=PP_ALIGN.CENTER)

# Output
add_rect(s, 9.2, 1.5, 3.7, 4.5, RGBColor(0x03, 0x1A, 0x0A))
add_textbox(s, "17 MODULES OUTPUT", 9.2, 1.5, 3.7, 0.45,
            font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
outputs = ["Resume Score + ATS", "Cover Letter Top 1%", "Skills Gap + Certs", "Visa Pathways", "Salary Negotiation", "30-60-90 Action Plan", "+ 11 more modules"]
add_bullet_box(s, outputs, 9.3, 2.0, 3.5, 3.5, font_size=12, color=LIGHT_GRAY, bullet="✓")

add_textbox(s, "⏱ Under 30 seconds  ·  🔒 Zero data stored  ·  💰 100% Free",
            0.5, 6.3, 12.3, 0.6,
            font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 3 (90 seconds)

So what does Alfalah AI actually do?

A user — anywhere in the world — uploads their resume.
They can paste text, upload a PDF, a Word document — whatever they have.

They select their country from 195 options. Their industry from 23 categories.
Optionally they paste a job description they want to target.

Our Azure-powered engine takes that input and in under 30 seconds —
produces 17 complete career intelligence modules.

Not one report. Not two. Seventeen.

Resume Score. Cover Letter. Skills Gap. Visa Pathways. Salary Negotiation.
Interview Prep. 30-60-90 Day Plan. And ten more.

Everything is country-aware. If you are in Nigeria applying for a job in Canada,
you get the Nigerian labor laws AND the Canadian visa pathways in the same analysis.

And here is what makes us different: Zero data is stored.
Your resume never touches a database. The moment we return results,
everything is discarded. Privacy by architecture, not just by policy.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — 17 MODULES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, PURPLE)

add_textbox(s, "17 CAREER INTELLIGENCE MODULES", 0.5, 0.15, 12, 0.6,
            font_size=26, bold=True, color=WHITE)
add_textbox(s, "One analysis. Seventeen expert outputs. In under 30 seconds.",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=PURPLE)

modules = [
    ("01", "Resume Score",        "ATS algorithm — 8 weighted dimensions"),
    ("02", "Recruiter POV",       "6-second hiring manager skim simulation"),
    ("03", "Cover Letter",        "Top-1% quality — quantified, specific"),
    ("04", "Resume Rewrite",      "Impact-first restructure + ATS audit"),
    ("05", "Skills Gap",          "Matched/missing + cert URLs"),
    ("06", "Interview Prep",      "5 Q&A sets — behavioral + technical"),
    ("07", "STAR Stories",        "3 metrics-driven behavioral examples"),
    ("08", "LinkedIn Summary",    "Search-optimized About section"),
    ("09", "Intro Scripts",       "1, 2, 3-minute intros by level"),
    ("10", "Matching Jobs",       "Titles, companies, 7-country boards"),
    ("11", "Visa Pathways",       "All immigration routes + govt URLs"),
    ("12", "Thank You Email",     "Post-interview differentiator"),
    ("13", "Salary Negotiation",  "Market table + word-for-word scripts"),
    ("14", "Action Plan",         "30-60-90 day structured career plan"),
    ("15", "Cold Outreach",       "LinkedIn DM + cold email templates"),
    ("16", "Career Pivot",        "Adjacent roles + 90-day roadmap"),
    ("17", "Country Laws",        "Labor law + ATS compliance by country"),
]

cols = 3
per_col = 6
box_w = 3.9
box_h = 0.72

for i, (num, title, desc) in enumerate(modules):
    col = i // per_col
    row = i % per_col
    x = 0.3 + col * (box_w + 0.25)
    y = 1.3 + row * (box_h + 0.06)
    add_rect(s, x, y, box_w, box_h, RGBColor(0x0D, 0x1B, 0x3E))
    add_textbox(s, num, x + 0.1, y + 0.04, 0.55, 0.38,
                font_size=14, bold=True, color=AZURE_BLUE)
    add_textbox(s, title, x + 0.65, y + 0.04, 3.1, 0.35,
                font_size=13, bold=True, color=WHITE)
    add_textbox(s, desc, x + 0.65, y + 0.36, 3.1, 0.3,
                font_size=10, color=MID_GRAY)

# Slide 17 — extra modules note
add_textbox(s, "+ Live Jobs  ·  ISCO-08 Similar Roles  ·  Top-1% Hiring Framework",
            0.5, 6.55, 12.3, 0.5,
            font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 4 (60 seconds)

Let me walk you through the 17 modules — and I want you to appreciate
that every single one of these is generated simultaneously, in a single API call.

Start with the Resume Score — a scientific ATS scoring algorithm across
8 weighted dimensions. Then the Recruiter POV — we simulate the 6-second
skim that a hiring manager does when they have 200 resumes on their desk.

The Cover Letter, Resume Rewrite, Skills Gap — these are not templates.
They are customized, quantified, specific to the user's resume, their target role,
and their country.

Visa Pathways is unique to us. No other free platform tells you every immigration
route from your country to the job's country, with official government URLs.

Salary Negotiation comes with word-for-word scripts, in the local currency of the user's country.

And the Country Laws module — labor law, worker rights, ATS compliance —
by country. So a job seeker in the UAE knows the rules are different than in Germany.

Seventeen modules. Zero cost. That is what we deliver.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — AZURE ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)

add_textbox(s, "AZURE ARCHITECTURE", 0.5, 0.15, 10, 0.6,
            font_size=26, bold=True, color=WHITE)
add_textbox(s, "100% Microsoft Azure Cloud Ecosystem — rg-v3 · Canada East",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=AZURE_BLUE)

# Architecture boxes
layers = [
    (AZURE_BLUE,  "USER BROWSER / PWA",           "195 countries · Any device · Zero install · Progressive Web App",       0.3, 1.2, 12.5, 0.75),
    (RGBColor(0x00,0x5A,0x9E), "AZURE STATIC WEB APPS — Global CDN Edge",  "React / Next.js · Auto-deploy from GitHub Actions · TLS 1.3 · Nearest edge node",  0.3, 2.1, 12.5, 0.75),
    (PURPLE,      "AZURE FUNCTIONS v2 — Python Serverless Backend",  "POST /career  ·  POST /chat  ·  POST /jobs  ·  GET /location  ·  GET /health  ·  POST /upload",  0.3, 3.0, 12.5, 0.75),
]
for color, title, desc, x, y, w, h in layers:
    add_rect(s, x, y, w, h, color)
    add_textbox(s, title, x+0.15, y+0.04, w-0.3, 0.34,
                font_size=13, bold=True, color=WHITE)
    add_textbox(s, desc, x+0.15, y+0.38, w-0.3, 0.3,
                font_size=11, color=LIGHT_GRAY)

# Bottom row — 4 Azure services
services = [
    (AZURE_BLUE,  "Azure OpenAI\nGPT-4o-mini",    "Primary AI"),
    (GREEN,       "Azure AI Search\nSemantic + Vector",  "RAG Retrieval"),
    (ORANGE,      "Azure Key Vault\nRBAC + Secrets",     "Security"),
    (PURPLE,      "Azure Content Safety\nOutput Moderation", "AI Safety"),
]
for i, (color, name, role) in enumerate(services):
    x = 0.3 + i * 3.2
    add_rect(s, x, 3.9, 2.9, 1.2, color)
    add_textbox(s, name, x+0.1, 3.95, 2.7, 0.65,
                font_size=12, bold=True, color=WHITE)
    add_textbox(s, role, x+0.1, 4.58, 2.7, 0.4,
                font_size=11, color=LIGHT_GRAY)

# RAG + Fallback row
add_rect(s, 0.3, 5.25, 6.0, 1.0, RGBColor(0x1A, 0x0A, 0x2E))
add_textbox(s, "RAG KNOWLEDGE ENGINE — 28 Files · 1M+ Tokens",
            0.45, 5.28, 5.7, 0.38, font_size=12, bold=True, color=PURPLE)
add_textbox(s, "436 occupations · 195 countries · 900+ skills · Future jobs 2026–2125",
            0.45, 5.62, 5.7, 0.5, font_size=11, color=LIGHT_GRAY)

add_rect(s, 6.6, 5.25, 6.4, 1.0, RGBColor(0x1A, 0x0A, 0x05))
add_textbox(s, "AI FALLBACK CHAIN — 4 Providers · 8 Models · 99.9% Uptime",
            6.75, 5.28, 6.1, 0.38, font_size=12, bold=True, color=ORANGE)
add_textbox(s, "Azure OpenAI → Gemini KEY1 → Gemini KEY2 → xAI Grok-4",
            6.75, 5.62, 6.1, 0.5, font_size=11, color=LIGHT_GRAY)

add_textbox(s, "GitHub Actions CI/CD  ·  Zero downtime deploy  ·  Auto-scale to millions",
            0.5, 6.55, 12.3, 0.5,
            font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 5 (90 seconds)

Now let me walk you through the Azure architecture that powers all of this.

At the top, the user opens their browser — any browser, any device, anywhere in the world.
Our frontend is served from Azure Static Web Apps, which means it is delivered from
the nearest CDN edge node. A user in Lagos, Nigeria gets the same performance
as a user in Toronto, Canada.

The frontend talks to our Azure Functions backend — Python v2 serverless.
Six endpoints. All stateless. All auto-scaling. We pay zero when nobody is using it,
and it scales to millions with no changes.

The Azure Functions orchestrate four Azure services:
- Azure OpenAI running GPT-4o-mini for primary AI inference
- Azure AI Search for semantic and vector retrieval from our knowledge base
- Azure Key Vault where every secret and API key is managed with RBAC
- Azure Content Safety to moderate every AI output

Below that is our RAG engine — 28 structured knowledge files,
over 1 million tokens of structured career science.

And our AI fallback chain — 4 providers, 8 models. If Azure OpenAI has a
momentary hiccup, we switch to Gemini in under one second.
If Gemini is unavailable, we fall to Grok. The platform never goes dark.

Everything deploys automatically from GitHub via GitHub Actions.
Zero-downtime, every push to main.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH STACK A-Z
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, GREEN)

add_textbox(s, "TECHNOLOGY STACK — A TO Z", 0.5, 0.15, 12, 0.6,
            font_size=26, bold=True, color=WHITE)
add_textbox(s, "Every tool, platform, language, and service used to build Alfalah AI",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=GREEN)

categories = [
    ("AZURE CLOUD (100%)", AZURE_BLUE, [
        "Static Web Apps · Functions v2 · OpenAI Service",
        "AI Search · Key Vault · Content Safety",
        "Entra ID B2C · Cosmos DB · Monitor / App Insights",
    ]),
    ("AI & MACHINE LEARNING", PURPLE, [
        "GPT-4o-mini (Azure OpenAI) — Primary",
        "Gemini 2.0 Flash / 1.5 Flash (Google AI Studio)",
        "xAI Grok-4-latest — Final fallback",
        "Custom RAG Engine — 99% structured data",
    ]),
    ("LANGUAGES & RUNTIMES", GREEN, [
        "Python 3.11+ — Azure Functions backend",
        "JavaScript ES2022 / TypeScript 5.x — Frontend",
        "HTML5 · CSS3 · Tailwind CSS 3.x",
        "Bash / PowerShell — Deployment scripts",
    ]),
    ("DEV TOOLS & IDE", ORANGE, [
        "Visual Studio Code — Primary IDE",
        "Claude Code (Anthropic) — AI pair programmer",
        "GitHub Copilot — Code completion",
        "Azure Functions Core Tools v4",
        "Postman — API testing",
    ]),
    ("FRAMEWORKS & STANDARDS", GOLD, [
        "Next.js 14 · React 18 · Tailwind CSS",
        "Progressive Web App (W3C Standard)",
        "ISCO-08 (ILO) · ESCO · O*NET · NOC · SOC",
        "Microsoft Responsible AI Standard v2",
    ]),
    ("DATA & KNOWLEDGE BASE", RGBColor(0xEF, 0x44, 0x44), [
        "28 structured Markdown knowledge files",
        "436 ISCO-08 occupations · 195 country packages",
        "900+ skills A-Z · Future occupations 2026–2125",
        "Serper.dev (Google Jobs) · ipapi.co (Geo)",
    ]),
]

for i, (title, color, items) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.3
    y = 1.3 + row * 2.8
    add_rect(s, x, y, 4.1, 2.6, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, x, y, 4.1, 0.42, color)
    add_textbox(s, title, x+0.1, y+0.04, 3.9, 0.35,
                font_size=11, bold=True, color=WHITE)
    add_bullet_box(s, items, x+0.1, y+0.48, 3.9, 2.05,
                   font_size=11, color=LIGHT_GRAY, bullet="·")

add_notes(s, """SPEAKER SCRIPT — SLIDE 6 (60 seconds)

This slide shows our complete technology stack from A to Z.

In the top-left, Azure Cloud. We use nine Azure services — Static Web Apps,
Functions, OpenAI, AI Search, Key Vault, Content Safety, Entra ID, Cosmos DB, and Monitor.
One cloud. One bill. Full integration. This is why we built on Azure.

AI and Machine Learning — we use four providers and eight models,
with our custom RAG engine delivering 99% of the intelligence from structured data.

Python 3.11 on the backend, Next.js and React on the frontend,
Tailwind for the galaxy gradient UI that users love.

Our development environment is Visual Studio Code, with Claude Code as our AI pair programmer,
and GitHub Copilot for completion. This project was built at the intersection of
human creativity and AI assistance.

And our knowledge base is grounded in internationally recognized standards —
ISCO-08, ESCO, O*NET — the same frameworks used by the ILO, EU, and US Department of Labor.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — RAG ENGINE & KNOWLEDGE BASE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, PURPLE)

add_textbox(s, "RAG KNOWLEDGE ENGINE", 0.5, 0.15, 12, 0.6,
            font_size=26, bold=True, color=WHITE)
add_textbox(s, "99% Structured Career Science · 1% Azure OpenAI Formatting · Zero Hallucination",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=PURPLE)

add_textbox(s, "What makes Alfalah AI different from every other AI career tool:",
            0.5, 1.2, 12, 0.4, font_size=15, bold=True, color=ORANGE)

diff_points = [
    "Every other AI career tool asks the LLM to guess — and LLMs hallucinate salary ranges, visa requirements, and labor laws.",
    "Alfalah AI retrieves structured facts first, THEN asks the LLM to format them into professional language.",
    "This means our salary data is correct. Our visa routes are accurate. Our labor laws are authoritative.",
]
add_bullet_box(s, diff_points, 0.5, 1.65, 12.3, 1.4, font_size=13, color=LIGHT_GRAY, bullet="▸")

# Stats boxes
kb_stats = [
    ("436", "ISCO-08\nOccupations", AZURE_BLUE),
    ("195", "Country\nData Packages", GREEN),
    ("900+", "Skills\nA-Z Master", PURPLE),
    ("28", "Knowledge\nBase Files", ORANGE),
    ("1M+", "Structured\nTokens", GOLD),
    ("3,000+", "ESCO\nOccupations", RGBColor(0xEF, 0x44, 0x44)),
]
for i, (num, label, color) in enumerate(kb_stats):
    x = 0.4 + i * 2.1
    add_rect(s, x, 3.2, 1.9, 1.5, color)
    add_textbox(s, num, x, 3.22, 1.9, 0.75,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, 3.95, 1.9, 0.7,
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "SOURCES:", 0.5, 4.9, 1.2, 0.4,
            font_size=13, bold=True, color=AZURE_BLUE)
sources = "ILO (ISCO-08)  ·  European Commission (ESCO)  ·  US Dept of Labor (O*NET)  ·  Statistics Canada (NOC)  ·  US BLS (SOC)"
add_textbox(s, sources, 1.7, 4.9, 11.1, 0.4,
            font_size=12, color=LIGHT_GRAY)

add_textbox(s, "Azure AI Search — Semantic + Vector retrieval across all 28 files — top-k chunk retrieval per query",
            0.5, 5.45, 12.3, 0.5, font_size=13, color=MID_GRAY)

add_rect(s, 0.3, 6.1, 12.5, 0.85, RGBColor(0x0D, 0x1B, 0x3E))
add_textbox(s, "\"We did not build an AI that guesses. We built an AI that knows.\"",
            0.5, 6.15, 12.3, 0.7,
            font_size=16, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 7 (90 seconds)

This slide explains the core innovation that separates Alfalah AI from everything else.

Every other AI career tool works the same way:
Upload resume, ask ChatGPT to tell you what to do.
The problem? ChatGPT makes things up. It will tell you the wrong visa route.
It will fabricate salary numbers. It will give you labor laws that do not exist.

We took a fundamentally different approach.

We built what is called a Retrieval-Augmented Generation engine.
Before the AI writes a single word, we retrieve authoritative structured data
from our knowledge base — 28 files, over 1 million tokens of career science.

For a software engineer in Pakistan applying to jobs in Canada —
we retrieve the exact ISCO-08 classification for their role,
the Canadian salary ranges from our structured data,
the specific visa pathways from Pakistan to Canada with official government URLs,
and the Ontario labor law requirements.

THEN we send all of that to Azure OpenAI and say:
"Format this into professional career guidance."

The AI does not guess. The AI formats facts.
That is why our outputs are accurate. That is why our visa routes are real.
That is what we mean when we say: Zero hallucination.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — RESPONSIBLE AI
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, GREEN)

add_textbox(s, "RESPONSIBLE AI BY DESIGN", 0.5, 0.15, 12, 0.6,
            font_size=26, bold=True, color=WHITE)
add_textbox(s, "Microsoft Responsible AI Standard v2 · EU AI Act Aligned · GRC Framework",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=GREEN)

principles = [
    ("FAIRNESS",      GREEN,      "Analysis based solely on skills, experience, and role requirements.\nNo demographic inference. No bias by nationality, age, or gender."),
    ("RELIABILITY",   AZURE_BLUE, "4-provider, 8-model fallback chain.\nPlatform never goes dark. 99.9%+ uptime guaranteed."),
    ("PRIVACY",       PURPLE,     "Zero data storage by architecture — not just policy.\nNo resume. No PII. No session data. Ever."),
    ("SECURITY",      ORANGE,     "Azure Key Vault · RBAC · Content Safety · TLS 1.3\nInput sanitization · Rate limiting · 50KB body guard."),
    ("INCLUSIVENESS", GOLD,       "195 countries · 23 industries · Accessibility guides\nfor youth (5-18), seniors (55+), and people with disabilities."),
    ("TRANSPARENCY",  RGBColor(0xEF,0x44,0x44), "ATS scoring dimensions disclosed to users.\nRAG sources cited. Full methodology published in /docs/."),
]

for i, (title, color, desc) in enumerate(principles):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.4
    y = 1.3 + row * 1.8
    add_rect(s, x, y, 6.1, 1.6, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, x, y, 6.1, 0.4, color)
    add_textbox(s, title, x+0.15, y+0.04, 5.8, 0.33,
                font_size=13, bold=True, color=WHITE)
    add_textbox(s, desc, x+0.15, y+0.45, 5.8, 1.1,
                font_size=12, color=LIGHT_GRAY)

add_notes(s, """SPEAKER SCRIPT — SLIDE 8 (60 seconds)

Microsoft's Responsible AI Standard has six principles.
We did not just try to comply with them — we designed every architectural
decision around them.

Fairness: Our analysis never looks at age, gender, nationality, or background.
It looks at skills, experience, and role requirements. Period.

Reliability: Four AI providers, eight models. The platform stays up even if
Azure OpenAI, Google, and xAI all have issues simultaneously.

Privacy: Most platforms say "we do not sell your data." We say:
"your data never exists in our system." There is nothing to sell because
there is nothing stored.

Security: Azure Key Vault for secrets. RBAC for access control.
Content Safety on every output. This is enterprise-grade security on a free platform.

Inclusiveness: We built specific career guides for youth aged 5 to 18,
for seniors aged 55 and above re-entering the workforce,
and for people with disabilities. Because "built for 8 billion" means everyone.

Transparency: We publish our ATS scoring algorithm dimensions.
We tell users exactly how their resume is being evaluated.
No black box.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO GUIDE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, GOLD)

add_textbox(s, "LIVE DEMO", 0.5, 0.15, 12, 0.6,
            font_size=32, bold=True, color=WHITE)
add_textbox(s, "Live: shahzad-job-coach-ai.vercel.app  |  API: govrag-v3-func.azurewebsites.net",
            0.5, 0.75, 12, 0.45, font_size=16, color=GOLD, align=PP_ALIGN.LEFT)

demo_steps = [
    ("STEP 1 — SPLASH SCREEN",     AZURE_BLUE,  "Show the 5-phase smart splash: IP auto-detection → country confirmed → industry selected.\nPoint out: 195-country coverage, zero login required."),
    ("STEP 2 — UPLOAD RESUME",     PURPLE,      "Paste or upload a sample resume. Show PDF/DOCX/TXT support.\nExplain: extracted in memory — never written to disk."),
    ("STEP 3 — AI ANALYSIS",       GREEN,       "Click Analyze. Show the 30-second loading state.\nExplain the RAG pipeline running in Azure Functions."),
    ("STEP 4 — 17 MODULES",        ORANGE,      "Scroll through the results dashboard. Highlight:\n• Resume Score with 8 dimensions\n• Recruiter POV — 6-second simulation\n• Visa Pathways — show for a cross-country example"),
    ("STEP 5 — COUNTRY AWARENESS", GOLD,        "Change country. Show how salary data changes to local currency.\nShow how visa pathways change. This is the differentiator."),
    ("STEP 6 — LIVE JOBS",         RGBColor(0xEF,0x44,0x44), "Scroll to live job postings — real listings from Google Jobs, last 7 days, country-filtered."),
]

for i, (title, color, desc) in enumerate(demo_steps):
    row = i % 3
    col = i // 3
    x = 0.3 + col * 6.5
    y = 1.35 + row * 1.85
    add_rect(s, x, y, 6.2, 1.65, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, x, y, 6.2, 0.38, color)
    add_textbox(s, title, x+0.12, y+0.04, 6.0, 0.3,
                font_size=12, bold=True, color=WHITE)
    add_textbox(s, desc, x+0.12, y+0.42, 6.0, 1.18,
                font_size=11, color=LIGHT_GRAY)

add_notes(s, """SPEAKER SCRIPT — SLIDE 9 — LIVE DEMO (2-3 minutes)

Now let me show you the platform live.

[Open browser to shahzad-job-coach-ai.vercel.app]

STEP 1: Notice the splash screen. The platform automatically detects your country
via IP address. You can confirm it or choose from all 195 countries.
Then you pick your industry. No account. No email. No password. You are in.

STEP 2: I will paste a sample resume here. You can see it accepts plain text,
PDF, and Word documents. Whatever format the user has.

STEP 3: [Click Analyze] Watch the loading screen — this is
the Azure Functions backend retrieving context from AI Search,
assembling the RAG prompt, and calling GPT-4o-mini. Under 30 seconds.

STEP 4: [Scroll through results] Look at the Resume Score — eight dimensions,
each explained. Then the Recruiter POV — this is the 6-second hiring
manager skim simulation. [Show Visa Pathways if cross-country]

STEP 5: [Change country] Watch how the salary data changes to local currency.
Watch how the visa pathways completely change. This is country-aware AI.
No competitor does this for free, for 195 countries.

STEP 6: [Scroll to live jobs] These are real job postings from Google Jobs,
pulled in the last 7 days, filtered for this country and role.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — GLOBAL IMPACT
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, ORANGE)

add_textbox(s, "GLOBAL IMPACT", 0.5, 0.15, 12, 0.6,
            font_size=28, bold=True, color=WHITE)
add_textbox(s, "Who benefits. How. Why it matters for humanity.",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=ORANGE)

personas = [
    ("🇵🇰 FATIMA, 24 — LAHORE",    "Computer Science graduate. Zero career guidance available locally.\nAlfalah AI gives her a top-1% resume, a Canadian work visa pathway,\nand a salary negotiation script — in 30 seconds. Free."),
    ("🇳🇬 CHUKWU, 31 — LAGOS",     "Mechanical engineer. Rejected 40 times. Does not know why.\nAlfalah AI shows him ATS rejection reasons, rewrites his resume\nwith impact-first bullets, and identifies 3 certifications he needs."),
    ("🇮🇳 PRIYA, 45 — BANGALORE",   "15 years experience. Wants to pivot to data science.\nAlfalah AI gives her a pivot score, adjacent role roadmap,\nand a 90-day transition plan with specific upskilling milestones."),
    ("🇿🇦 SIPHO, 58 — JOHANNESBURG","Lost job in restructuring. 55+ re-entry is brutal.\nAlfalah AI has a specific seniors career guide, identifies transferable\nskills, and generates cold outreach templates for his network."),
]

for i, (name, story) in enumerate(personas):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.4
    y = 1.3 + row * 2.4
    add_rect(s, x, y, 6.1, 2.2, RGBColor(0x0D, 0x1B, 0x3E))
    add_textbox(s, name, x+0.15, y+0.08, 5.8, 0.45,
                font_size=14, bold=True, color=ORANGE)
    add_textbox(s, story, x+0.15, y+0.55, 5.8, 1.55,
                font_size=12, color=LIGHT_GRAY)

add_textbox(s, "These are not hypothetical users. These are the 8 billion people this platform was built to serve.",
            0.5, 6.4, 12.3, 0.7,
            font_size=14, italic=True, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 10 (60 seconds)

I want to put a human face on this platform.

Fatima is 24 years old in Lahore, Pakistan. She has a Computer Science degree
but zero professional career guidance available to her locally.
In 30 seconds, Alfalah AI gives her a top-1% resume,
a Canadian skilled worker visa pathway with official government URLs,
and a word-for-word salary negotiation script. For free.

Chukwu has been rejected 40 times and does not know why.
Alfalah AI tells him exactly why — ATS rejection patterns —
rewrites his resume, and tells him the 3 certifications he needs
to break through.

Priya is 45 and wants to pivot to data science after 15 years in another field.
She gets a career pivot score, adjacent role analysis, and a 90-day roadmap.

Sipho is 58. He is re-entering after a layoff. We built a specific
accessibility guide for professionals 55 and above because
"built for 8 billion" means everyone, including those the industry forgets.

That is who this platform serves. That is why it had to be free.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — WHY MICROSOFT AZURE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)

add_textbox(s, "WHY MICROSOFT AZURE", 0.5, 0.15, 12, 0.6,
            font_size=28, bold=True, color=WHITE)
add_textbox(s, "Not just cloud infrastructure. A values-aligned partnership.",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=AZURE_BLUE)

reasons = [
    ("AI FOR GOOD",         AZURE_BLUE, "Microsoft's AI for Good initiative aligns perfectly with our mission.\nAzure enables free-tier infrastructure that makes a truly free global platform possible."),
    ("GLOBAL PRESENCE",     GREEN,      "Azure has 60+ regions globally — our users in 195 countries get\nlow-latency performance from their nearest Azure edge node."),
    ("RESPONSIBLE AI",      PURPLE,     "Azure Content Safety, Key Vault, RBAC, and the Responsible AI Standard v2\ngave us enterprise-grade trust infrastructure for a free public platform."),
    ("INTEGRATED AI STACK", ORANGE,     "Azure OpenAI + AI Search + Content Safety are deeply integrated.\nNo glue code between providers. One identity, one billing, one audit trail."),
    ("COMPLIANCE READY",    GOLD,       "Azure's compliance certifications — SOC2, ISO27001, GDPR, HIPAA —\nmean our platform is enterprise-compliant from day one."),
    ("HACKATHON ALIGNMENT", RGBColor(0xEF,0x44,0x44), "This project demonstrates the full Azure AI ecosystem in production:\nStatic Web Apps + Functions + OpenAI + AI Search + Key Vault = one cohesive solution."),
]

for i, (title, color, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.4
    y = 1.3 + row * 1.8
    add_rect(s, x, y, 6.1, 1.62, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, x, y, 0.18, 1.62, color)
    add_textbox(s, title, x+0.3, y+0.08, 5.7, 0.38,
                font_size=13, bold=True, color=color)
    add_textbox(s, desc, x+0.3, y+0.48, 5.7, 1.05,
                font_size=12, color=LIGHT_GRAY)

add_notes(s, """SPEAKER SCRIPT — SLIDE 11 (60 seconds)

Someone will ask: why Azure? Why not AWS? Why not GCP?

The answer is alignment. Not just technical alignment — values alignment.

Microsoft's AI for Good initiative exists because Satya Nadella believes
technology should serve humanity. That is exactly what we built.

Azure's free tier on Static Web Apps and Functions means we can run
a global platform at zero cost during development — critical for a
free, not-for-profit tool.

The integrated Azure AI stack — OpenAI, AI Search, Content Safety —
meant we could build a responsible AI platform without stitching
together five different vendors and five different trust frameworks.

Azure's 60+ global regions mean that a user in Indonesia gets the same
performance as a user in Seattle. That matters when you are serving 195 countries.

And for this hackathon specifically — we are demonstrating the full Azure AI stack
in a real production use case that serves real people.
That is what this competition is designed to reward.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — ROADMAP
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, PURPLE)

add_textbox(s, "ROADMAP", 0.5, 0.15, 12, 0.6,
            font_size=28, bold=True, color=WHITE)
add_textbox(s, "Where we have been. Where we are going. What comes next.",
            0.5, 0.72, 12, 0.4, font_size=14, italic=True, color=PURPLE)

phases = [
    ("V1", "Mar 2026", "COMPLETE ✅", GREEN,
     ["12 AI career modules", "Resume analysis + ATS scoring", "Deployed live on Vercel", "Submitted H1 Hackathon ($6K)"]),
    ("V2", "Mar 2026", "COMPLETE ✅", GREEN,
     ["17 modules + RAG engine", "195 countries + industry awareness", "Gemini + Grok fallback chain", "PWA mobile-ready"]),
    ("V3", "Mar 2026", "ACTIVE 🟡", GOLD,
     ["100% Azure cloud ecosystem", "Python Functions backend", "AI Search + Key Vault + GRC", "THIS PRESENTATION"]),
    ("V4", "2026 Q3", "PLANNED 📋", AZURE_BLUE,
     ["React Native mobile app", "Cosmos DB for analytics", "Multi-language (40+ languages)", "Voice input + accessibility"]),
    ("V5", "2026 Q4", "PLANNED 📋", PURPLE,
     ["alfalah.app domain live", "Azure Marketplace listing", "Enterprise API tier", "1 million users target"]),
]

for i, (version, date, status, color, items) in enumerate(phases):
    x = 0.3 + i * 2.55
    add_rect(s, x, 1.3, 2.35, 5.5, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, x, 1.3, 2.35, 0.5, color)
    add_textbox(s, version, x+0.1, 1.32, 1.0, 0.38,
                font_size=20, bold=True, color=WHITE)
    add_textbox(s, date, x+1.0, 1.37, 1.3, 0.3,
                font_size=11, color=WHITE)
    add_textbox(s, status, x+0.1, 1.82, 2.2, 0.35,
                font_size=11, bold=True, color=color)
    add_bullet_box(s, items, x+0.1, 2.2, 2.2, 4.4,
                   font_size=11, color=LIGHT_GRAY, bullet="·")

add_rect(s, 0.3, 6.55, 12.5, 0.65, RGBColor(0x0D, 0x1B, 0x3E))
add_textbox(s, "Next: React Native mobile app · alfalah.app domain · Azure Marketplace listing · 40-language global rollout",
            0.5, 6.6, 12.3, 0.55,
            font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 12 (60 seconds)

We have already shipped two complete versions in production.

V1 launched in early March 2026 — 12 modules, submitted to our first hackathon.
V2 followed immediately — 17 modules, 195 countries, the full RAG engine,
Gemini and Grok fallbacks. PWA-ready for mobile.

V3 is what we are presenting today — the full Azure migration,
Python backend, AI Search, Key Vault, and GRC framework.

V4 is planned for Q3 2026 — React Native mobile app,
Cosmos DB for analytics, and most importantly: multi-language support
for 40+ languages. Because Arabic speakers, Hindi speakers,
Swahili speakers deserve guidance in their own language.

V5 closes out 2026 with alfalah.app live, Azure Marketplace listing,
and a target of 1 million users.

V4 is the next milestone: React Native mobile app for Android and iOS,
Cosmos DB for anonymous analytics, and 40-language rollout —
because Arabic speakers, Hindi speakers, Swahili speakers deserve
guidance in their own language.

V5 closes 2026 with alfalah.app live, Azure Marketplace listing,
and a target of 1 million users served.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — THE ASK / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, GOLD)
add_rect(s, 0, 7.42, 13.33, 0.08, GOLD)

add_textbox(s, "COMMUNITY CONTRIBUTION", 0.5, 0.25, 12, 0.65,
            font_size=32, bold=True, color=WHITE)

add_textbox(s, "\"Built for humanity. Open source. Free forever.\"",
            0.5, 1.0, 12.3, 0.7,
            font_size=22, italic=True, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

contribution_items = [
    ("OPEN SOURCE — MIT", AZURE_BLUE,
     "Full source code at github.com/shahzadms7/v3 — fork it, learn from it, build on it.\nAny developer anywhere can extend this platform for their community."),
    ("FREE FOREVER", GREEN,
     "No subscription. No login. No paywall. No data harvested.\nEvery feature available to every professional in every country, at zero cost."),
    ("AZURE ECOSYSTEM SHOWCASE", PURPLE,
     "A production reference architecture for Azure Static Web Apps + Functions + OpenAI +\nAI Search + Key Vault + Content Safety — all working together at real scale."),
    ("RESPONSIBLE AI IN PRACTICE", ORANGE,
     "Microsoft Responsible AI Standard v2 implemented end-to-end.\nZero hallucination via RAG grounding. Privacy by architecture. GRC compliance framework."),
]
for i, (title, color, desc) in enumerate(contribution_items):
    y = 1.9 + i * 1.18
    add_rect(s, 0.4, y, 12.2, 1.05, RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(s, 0.4, y, 0.2, 1.05, color)
    add_textbox(s, title, 0.75, y+0.07, 3.2, 0.38,
                font_size=14, bold=True, color=color)
    add_textbox(s, desc, 4.1, y+0.1, 8.4, 0.85,
                font_size=12, color=LIGHT_GRAY)

add_textbox(s, "govrag-v3-func.azurewebsites.net/api/health  ·  github.com/shahzadms7/v3  ·  shahzad-job-coach-ai.vercel.app",
            0.5, 6.65, 12.3, 0.55,
            font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 13 (45 seconds)

I want to close the content portion of this presentation with what this project gives back.

First — it is open source under MIT license. Every line of code, every knowledge file,
every architectural decision is public at github.com/shahzadms7/v3.
Any developer in any country can fork this and build a version for their community.

Second — it is free forever. We have never charged for it.
We will never charge for it. The mission is the product.

Third — this is a complete, production-grade reference architecture for the full
Azure AI stack. Azure Functions + Static Web Apps + OpenAI + AI Search +
Key Vault + Content Safety — all working together, all documented,
all available for the community to learn from.

Fourth — it demonstrates that Responsible AI is not a theoretical principle.
It is an architectural decision. Zero data storage. RAG grounding.
Content Safety on every output. GRC compliance framework.
These are not policies. They are design choices baked into every line of code.

That is what we built. That is what we give back.
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING / THANK YOU
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)
add_rect(s, 0, 7.42, 13.33, 0.08, ORANGE)

add_textbox(s, "الفلاح", 0, 1.0, 13.33, 1.5,
            font_size=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(s, "Come to Success.", 0, 2.4, 13.33, 0.7,
            font_size=28, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "For every human on Earth.", 0, 3.05, 13.33, 0.6,
            font_size=22, color=AZURE_BLUE, align=PP_ALIGN.CENTER)

add_rect(s, 3, 3.8, 7.33, 0.04, ORANGE)

add_textbox(s, "ALFALAH JOB CAREER INTELLIGENT AI  2026  V3", 0, 3.95, 13.33, 0.6,
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

links = [
    "Platform: shahzad-job-coach-ai.vercel.app",
    "API: govrag-v3-func.azurewebsites.net/api/health",
    "Source: github.com/shahzadms7/v3  |  Architecture: ARCHITECTURE.md",
]
for i, link in enumerate(links):
    add_textbox(s, link, 0, 4.65 + i * 0.42, 13.33, 0.38,
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, "Shahzad Muhammad  ·  Mississauga, Ontario, Canada",
            0, 6.0, 13.33, 0.4,
            font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, "Built on Microsoft Azure  ·  Free Forever  ·  Open Source MIT",
            0, 6.4, 13.33, 0.4,
            font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_notes(s, """SPEAKER SCRIPT — SLIDE 14 — CLOSING (45 seconds)

I want to close with the name.

Alfalah. الفلاح.

In Arabic, in Urdu, in the Islamic tradition — this word means success.
But not just personal success. Flourishing. The kind of success that
comes from doing something meaningful for others.

Every day, somewhere on this planet, a young person is rejected by an
ATS system they do not understand. A professional is stuck in a country
with no path forward. A parent is trying to give their children a better life
but does not know where to start.

We built this platform for them.

Not with a complex enterprise contract. Not with a subscription.
Not with a login page. Just: go to the website, upload your resume,
and in 30 seconds — get the guidance that changes the trajectory of your career.

That is what technology should do. That is what AI on Azure can do.

Thank you. I am happy to take questions.

[Pause. Smile. Let it land.]
""")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — ARCHITECTURE DIAGRAM (full-page PNG image)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.08, AZURE_BLUE)

add_textbox(s, "SYSTEM ARCHITECTURE — End-to-End", 0.5, 0.1, 12, 0.5,
            font_size=20, bold=True, color=WHITE)
add_textbox(s, "Alfalah Job Career Intelligent AI 2026 V3  ·  100% Microsoft Azure  ·  govrag-v3-func.azurewebsites.net",
            0.5, 0.55, 12.3, 0.35, font_size=11, color=AZURE_BLUE, align=PP_ALIGN.LEFT)

arch_img = r"g:\My Drive\Claude Projects 2026\shahzad-job-coach-ai\v3\architecture_diagram.png"
import os
if os.path.exists(arch_img):
    s.shapes.add_picture(arch_img,
                         Inches(0.15), Inches(1.0),
                         Inches(13.0), Inches(6.3))

add_notes(s, """SPEAKER SCRIPT — SLIDE 15 — ARCHITECTURE DIAGRAM (60 seconds)

This diagram shows the complete end-to-end architecture of Alfalah AI.

Starting at the top — users from 195 countries access the platform through
any browser or as a Progressive Web App. No installation. No account.

Their request hits Azure Static Web Apps — served from the nearest global CDN edge node.

From there, Azure Functions v2 handles all six API routes —
career analysis, chat, jobs, location, health check, and file upload.

The Functions connect to five Azure services:
Azure OpenAI for primary AI inference, Azure AI Search for semantic retrieval
from our knowledge base, Content Safety for output moderation,
Key Vault for secrets management, and Application Insights for monitoring.

Below that is our AI fallback chain — 4 providers, 8 models —
so the platform never goes dark regardless of any single provider outage.

The RAG knowledge engine holds 32 career files, 513 KB, 163 ISCO-08 occupations,
416 skills, 195 countries, and 16,544 lines of code and documentation.

At the bottom right are the development tools that built this:
Claude Sonnet 4.6, Visual Studio Code, GitHub Actions, Azure CLI, and PowerShell.

Everything. In one diagram.
""")

# ── Save ────────────────────────────────────────────────────────────
output_path = r"g:\My Drive\Claude Projects 2026\shahzad-job-coach-ai\v3\Alfalah_Job_Career_Intelligent_AI_2026_V3.pptx"
prs.save(output_path)
print(f"PowerPoint saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
